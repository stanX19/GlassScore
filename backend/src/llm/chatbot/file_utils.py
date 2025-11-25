"""
Utility functions for extracting text from base64-encoded files.
Supports: PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), and plain text files.
"""
import base64
import io
from typing import Optional

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_from_base64(base64_content: str, mime_type: str) -> str:
    """
    Extract text from a base64-encoded file.
    
    Args:
        base64_content: Base64-encoded file content
        mime_type: MIME type of the file (e.g., 'application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')
    
    Returns:
        Extracted text content
    
    Raises:
        ValueError: If file type is unsupported or required library is not installed
    """
    try:
        file_bytes = base64.b64decode(base64_content)
        file_stream = io.BytesIO(file_bytes)
    except Exception as e:
        raise ValueError(f"Failed to decode base64 content: {e}")
    
    # PDF files
    if mime_type == "application/pdf":
        return _extract_pdf(file_stream)
    
    # Word documents
    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        return _extract_docx(file_stream)
    
    # Excel files
    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel"
    ]:
        return _extract_xlsx(file_stream)
    
    # PowerPoint files
    elif mime_type in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]:
        return _extract_pptx(file_stream)
    
    # Plain text files
    elif mime_type.startswith("text/"):
        try:
            return file_bytes.decode('utf-8')
        except UnicodeDecodeError:
            # Try other common encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    return file_bytes.decode(encoding)
                except UnicodeDecodeError:
                    continue
            raise ValueError("Failed to decode text file with common encodings")
    
    else:
        raise ValueError(f"Unsupported MIME type: {mime_type}")


def _extract_pdf(file_stream: io.BytesIO) -> str:
    """Extract text from PDF file."""
    if pdfplumber is None:
        raise ValueError("pdfplumber is not installed. Run: pip install pdfplumber")
    
    text_parts = []
    with pdfplumber.open(file_stream) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
    
    if not text_parts:
        return "[No text content found in PDF]"
    
    return "\n\n".join(text_parts)


def _extract_docx(file_stream: io.BytesIO) -> str:
    """Extract text from Word document."""
    if Document is None:
        raise ValueError("python-docx is not installed. Run: pip install python-docx")
    
    doc = Document(file_stream)
    text_parts = []
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text)
    
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text_parts.append(row_text)
    
    if not text_parts:
        return "[No text content found in Word document]"
    
    return "\n\n".join(text_parts)


def _extract_xlsx(file_stream: io.BytesIO) -> str:
    """Extract text from Excel file."""
    if load_workbook is None:
        raise ValueError("openpyxl is not installed. Run: pip install openpyxl")
    
    workbook = load_workbook(file_stream, data_only=True)
    text_parts = []
    
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text_parts.append(f"=== Sheet: {sheet_name} ===")
        
        for row in sheet.iter_rows(values_only=True):
            # Filter out empty cells and convert to strings
            row_values = [str(cell) for cell in row if cell is not None and str(cell).strip()]
            if row_values:
                text_parts.append(" | ".join(row_values))
    
    if len(text_parts) <= len(workbook.sheetnames):  # Only sheet headers
        return "[No data found in Excel file]"
    
    return "\n".join(text_parts)


def _extract_pptx(file_stream: io.BytesIO) -> str:
    """Extract text from PowerPoint presentation."""
    if Presentation is None:
        raise ValueError("python-pptx is not installed. Run: pip install python-pptx")
    
    prs = Presentation(file_stream)
    text_parts = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_parts.append(f"=== Slide {i} ===")
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    
    if len(text_parts) == 0:
        return "[No text content found in PowerPoint]"
    
    return "\n\n".join(text_parts)


def get_supported_mime_types() -> list[str]:
    """Return list of supported MIME types."""
    return [
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "text/plain",
        "text/csv",
        "text/html",
        "text/markdown",
    ]
